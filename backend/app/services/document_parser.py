from dataclasses import dataclass
from pathlib import Path

from fastapi import HTTPException, status

from app.core.config import BASE_DIR


@dataclass
class ParsedSection:
    page: int
    text: str


@dataclass
class ParsedDocument:
    text: str
    page_count: int
    sections: list[ParsedSection]


@dataclass
class ExtractedText:
    path: Path
    page_count: int
    character_count: int
    preview: str
    sections: list[ParsedSection]


def parse_document(file_path: Path, file_type: str) -> ParsedDocument:
    if not file_path.exists():
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Document file was not found.",
        )

    parsers = {
        "pdf": _parse_pdf,
        "docx": _parse_docx,
        "pptx": _parse_pptx,
        "txt": _parse_plain_text,
        "md": _parse_plain_text,
    }
    parser = parsers.get(file_type)
    if parser is None:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Unsupported file type: {file_type}.",
        )

    sections = parser(file_path)
    text = _format_sections(sections)

    return ParsedDocument(
        text=text,
        page_count=len(sections) if sections else 0,
        sections=sections,
    )


def extract_text_from_document(
    document_id: str,
    file_path: Path,
    file_type: str,
) -> ExtractedText:
    if not _is_valid_document_id(document_id):
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Invalid document_id.",
        )

    parsed_document = parse_document(file_path, file_type)
    text_dir = BASE_DIR / "extracted_text"
    text_dir.mkdir(parents=True, exist_ok=True)
    text_path = text_dir / f"{document_id}.txt"
    text_path.write_text(parsed_document.text, encoding="utf-8")

    return ExtractedText(
        path=text_path,
        page_count=parsed_document.page_count,
        character_count=len(parsed_document.text),
        preview=parsed_document.text[:500],
        sections=parsed_document.sections,
    )


def extract_text_from_pdf(document_id: str) -> ExtractedText:
    return extract_text_from_document(
        document_id=document_id,
        file_path=BASE_DIR / "uploads" / f"{document_id}.pdf",
        file_type="pdf",
    )


def _parse_pdf(file_path: Path) -> list[ParsedSection]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="PDF parsing dependency is not installed. Run: pip install -r requirements.txt",
        ) from exc

    reader = PdfReader(str(file_path))
    return [
        ParsedSection(page=page_index, text=(page.extract_text() or "").strip())
        for page_index, page in enumerate(reader.pages, start=1)
    ]


def _parse_docx(file_path: Path) -> list[ParsedSection]:
    try:
        from docx import Document
    except ImportError as exc:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="DOCX parsing dependency is not installed. Run: pip install -r requirements.txt",
        ) from exc

    document = Document(str(file_path))
    parts = []

    parts.extend(
        paragraph.text.strip()
        for paragraph in document.paragraphs
        if paragraph.text.strip()
    )

    for table in document.tables:
        for row in table.rows:
            cells = [
                cell.text.strip()
                for cell in row.cells
                if cell.text.strip()
            ]
            if cells:
                parts.append(" | ".join(cells))

    return [ParsedSection(page=1, text="\n".join(parts).strip())]


def _parse_pptx(file_path: Path) -> list[ParsedSection]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="PPTX parsing dependency is not installed. Run: pip install -r requirements.txt",
        ) from exc

    presentation = Presentation(str(file_path))
    sections = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

        sections.append(
            ParsedSection(
                page=slide_index,
                text="\n".join(texts).strip(),
            )
        )

    return sections


def _parse_plain_text(file_path: Path) -> list[ParsedSection]:
    text = _read_text_file(file_path)
    return [ParsedSection(page=1, text=text.strip())]


def _read_text_file(file_path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gb18030"):
        try:
            return file_path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue

    return file_path.read_text(encoding="utf-8", errors="replace")


def _format_sections(sections: list[ParsedSection]) -> str:
    return "\n\n".join(
        f"--- Page {section.page} ---\n{section.text}"
        for section in sections
    ).strip()


def _is_valid_document_id(document_id: str) -> bool:
    return len(document_id) == 64 and all(char in "0123456789abcdef" for char in document_id)
